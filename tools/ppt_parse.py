import os
import json
import time
import requests
from pptx import Presentation

# ----------------------------
# API Key & Model Setup
# ----------------------------
GROQ_API_KEY = "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"  # or replace with your key directly
GROQ_MODEL = "llama3-70b-8192"

if not GROQ_API_KEY:
    raise ValueError("❌ GROQ_API_KEY is not set. Please export it as an environment variable.")

# ----------------------------
# STEP 1: Extract Raw Slide Text
# ----------------------------
def extract_raw_slide_text(pptx_path):
    prs = Presentation(pptx_path)
    raw_slides = []

    for i, slide in enumerate(prs.slides):
        text_blocks = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    text_blocks.append(text)
        raw_slides.append({
            "slide_number": i + 1,
            "raw_text": "\n".join(text_blocks)
        })

    return raw_slides

# ----------------------------
# STEP 2: Groq API for Structuring
# ----------------------------
def structure_slide_with_llm_groq(raw_slide_text):
    url = "https://api.groq.com/openai/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {GROQ_API_KEY}",
        "Content-Type": "application/json"
    }

    prompt = f"""
You are a helpful assistant.

Given the following PowerPoint slide content, return a clean, structured version of the text. 
Remove all extra line breaks, bullet points, formatting symbols, and unnecessary labels like 'Output:', 'Slide Content:', or headings. 
Return only the raw content in logical reading order, like handwritten lecture notes. 
Do not include any bullets, hyphens, line breaks, or your own commentary.

Slide Content:
{raw_slide_text}
"""

    data = {
        "model": GROQ_MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.3,
        "max_tokens": 512
    }

    try:
        response = requests.post(url, headers=headers, json=data)
        response.raise_for_status()
        return response.json()["choices"][0]["message"]["content"].strip()
    except requests.exceptions.RequestException as e:
        print(f"⚠️ Error from Groq API: {e}")
        return ""

# ----------------------------
# STEP 3: Post-processing Cleanup
# ----------------------------
def clean_response(text):
    text = text.replace("Output:", "").replace("\n", " ").strip()
    return " ".join(text.split())

# ----------------------------
# STEP 4: Main Function
# ----------------------------
def generate_structured_notes_groq(pptx_path, output_json="structured_notes.json"):
    slides = extract_raw_slide_text(pptx_path)
    final_notes = []

    for slide in slides:
        print(f"⚙️  Processing Slide {slide['slide_number']}...")
        structured = structure_slide_with_llm_groq(slide["raw_text"])
        structured_clean = clean_response(structured)
        time.sleep(1)  # gentle rate limit buffer
        final_notes.append({
            "slide_number": slide["slide_number"],
            "structured_content": structured_clean
        })

    with open(output_json, "w") as f:
        json.dump({"slides": final_notes}, f, indent=2)

    print(f"\n✅ Structured notes saved to {output_json}")

# ----------------------------
# Run Script
# ----------------------------
if __name__ == "__main__":
    ppt_path = "../assets/29_S-attributed_SDD - Syntax tree Generation.pptx"
    generate_structured_notes_groq(ppt_path)
