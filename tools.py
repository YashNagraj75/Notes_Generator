import json
import os
import subprocess
import time

import requests
from agents import function_tool
from pptx import Presentation
from groq import Groq
from google import genai
from google.genai import types
from googleapiclient.discovery import build
import base64
from urllib.parse import urlparse


GROQ_API_KEY = os.environ.get("GROQ_API_KEY")  # or replace with your key directly
GROQ_MODEL = "llama3-70b-8192"
client = Groq(api_key=os.environ.get("GROQ_API_KEY"))

if not GROQ_API_KEY:
    raise ValueError(
        "❌ GROQ_API_KEY is not set. Please export it as an environment variable."
    )


def extract_raw_slide_text(pptx_path):
    prs = Presentation(pptx_path)
    raw_slides = []

    for i, slide in enumerate(prs.slides):
        text_blocks = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    text_blocks.append(text)
        raw_slides.append({"slide_number": i + 1, "raw_text": "\n".join(text_blocks)})

    return raw_slides


def structure_slide_with_llm_groq(raw_slide_text):
    url = "https://api.groq.com/openai/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {GROQ_API_KEY}",
        "Content-Type": "application/json",
    }

    prompt = f"""
You are a helpful assistant.

Given the following PowerPoint slide content, return a clean, structured version of the text. 
Remove all extra line breaks, bullet points, formatting symbols, and unnecessary labels like 'Output:', 'Slide Content:', or headings. 
Return only the raw content in logical reading order, like handwritten lecture notes. 
Do not include any bullets, hyphens, line breaks, or your own commentary.

Slide Content:
{raw_slide_text}
"""

    data = {
        "model": GROQ_MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.3,
        "max_tokens": 512,
    }

    try:
        response = requests.post(url, headers=headers, json=data)
        response.raise_for_status()
        return response.json()["choices"][0]["message"]["content"].strip()
    except requests.exceptions.RequestException as e:
        print(f"⚠️ Error from Groq API: {e}")
        return ""


def clean_response(text):
    text = text.replace("Output:", "").replace("\n", " ").strip()
    return " ".join(text.split())


async def generate_structured_output(pptx_path: str):
    """
    Generate structured notes from a PowerPoint file.
        Args:
            pptx_path (str): Path to the PowerPoint file.
        Returns:
            list: List of structured notes for each slide.
    """
    slides = extract_raw_slide_text(pptx_path)
    final_notes = []

    for slide in slides:
        print(f"⚙️  Processing Slide {slide['slide_number']}...")
        structured = structure_slide_with_llm_groq(slide["raw_text"])
        structured_clean = clean_response(structured)
        time.sleep(1)  # gentle rate limit buffer
        final_notes.append(
            {
                "slide_number": slide["slide_number"],
                "structured_content": structured_clean,
            }
        )

    return final_notes


def extract_audio_ffmpeg(input_video, output_audio):
    command = [
        "ffmpeg",
        "-i",
        input_video,
        "-vn",  # Disable video recording
        "-acodec",
        "libmp3lame",  # Audio codec
        "-q:a",
        "0",  # Audio quality
        output_audio,
    ]
    subprocess.run(command, check=True)


async def get_video_transcript(path: str):
    """
    Generate a transcript from a video file.
    Args:
        path (str): Path to the video file.
    Returns:
        str: Transcription of the video.
    """
    extract_audio_ffmpeg(path, "output_audio.wav")
    filename = os.path.dirname(__file__) + "/output_audio.wav"
    with open(filename, "rb") as file:
        transcription = client.audio.transcriptions.create(
            file=(filename, file.read()),
            model="distil-whisper-large-v3-en",
            response_format="verbose_json",
        )
        print(transcription.text)
        return transcription.text


@function_tool
def get_image_description(query: str, num_images: int) -> list[dict]:
    """
    Fetch image URLs and descriptions from Google Custom Search API

    Args:
        query: str: search query
        num_images: int: number of images to fetch

    Returns:
        List[dict]: List of dicts with 'url' and 'description' keys
    """
    api_key = os.environ.get("GEMINI_API_KEY")
    cse_id = os.environ.get("CSE_ID")

    # Build the service object
    service = build("customsearch", "v1", developerKey=api_key)

    # Perform the search
    try:
        res = (
            service.cse()
            .list(q=query, cx=cse_id, searchType="image", num=num_images)
            .execute()
        )
    except Exception:
        return []

    image_desc = []
    if "items" in res:
        for item in res["items"]:
            url = item.get("link")
            if url:
                parsed_url = urlparse(url)
                if parsed_url.scheme in ["http", "https"]:
                    try:
                        response = requests.get(url)
                        if response.status_code != 200:
                            continue
                        image_base64 = base64.b64encode(response.content).decode(
                            "utf-8"
                        )
                        # Get description using get_description function
                        description = get_description(image_base64)
                        image_desc.append(
                            {
                                "url": url,
                                "description": description,
                            }
                        )
                    except requests.RequestException as e:
                        print(f"Failed to fetch {url}: {e}")
                else:
                    print(f"Unsupported URL scheme {parsed_url.scheme} for URL: {url}")
        print(image_desc)
        return image_desc
    else:
        return []


def get_description(img_content):
    image = types.Part.from_bytes(
        data=base64.b64decode(img_content),
        mime_type="image/png",
    )
    gem_client = genai.Client(
        api_key=os.environ.get("GEMINI_API_KEY"),
    )

    model = "gemini-2.0-flash-001"
    contents = [
        types.Content(
            role="user",
            parts=[
                types.Part.from_text(
                    text="Generate a concise description of this image in about 50 words.",
                ),
                image,
            ],
        )
    ]

    response = gem_client.models.generate_content(model=model, contents=contents)
    return response.text
